"""Tools for the OpalaTex Autonomous Agent (MemGPT pattern + HITL)."""
import os
import json
import re
import subprocess
import sys
import time
from dataclasses import dataclass, field
from fnmatch import fnmatch
from pathlib import Path

import functools
import asyncio
import concurrent.futures
from agenticblocks.core.function_block import as_tool

from . import terminal as T

# ─── Shared progress state ───────────────────────────────────────────────────
# The orchestrator reads these to render a live progress panel.

TURN_ACHIEVEMENTS = ""



class _AgentProgress:
    def __init__(self):
        self.heartbeat: int = 0
        self.max_heartbeats: int = 15
        self.tasks_done: int = 0
        self.tasks_total: int = 0
        self.last_tool: str = "—"
        self.last_args: str = ""
        self.start_time: float = time.monotonic()
        self.live_context = None
        # Error / failure tracking
        self.task_failures: int = 0          # reviewer said task NOT DONE
        self.worker_errors: int = 0          # worker returned [ERROR ...]
        self.lint_errors: int = 0            # lint check failures from write/edit
        self.plan_retries: int = 0           # oracle reflection retries
        self.recent_events: list[str] = []   # last N notable events (errors, retries)

    def _push_event(self, msg: str) -> None:
        self.recent_events.append(msg)
        if len(self.recent_events) > 6:
            self.recent_events.pop(0)

    def record_lint_error(self, path: str) -> None:
        self.lint_errors += 1
        short = os.path.basename(path)
        self._push_event(f"[red]✗ lint[/red] {short}")

    def record_worker_error(self, task_id: str, snippet: str = "") -> None:
        self.worker_errors += 1
        snip = snippet[:60].replace("\n", " ") if snippet else ""
        self._push_event(f"[red]✗ worker[/red] {task_id}" + (f": {snip}" if snip else ""))

    def record_task_failure(self, task_id: str, reason: str = "") -> None:
        self.task_failures += 1
        snip = reason[:60].replace("\n", " ") if reason else ""
        self._push_event(f"[yellow]⚠ review[/yellow] {task_id}" + (f": {snip}" if snip else ""))

    def record_plan_retry(self, attempt: int) -> None:
        self.plan_retries += 1
        self._push_event(f"[yellow]⚠ plan retry[/yellow] #{attempt}")

    def update(self, tool_name: str, args_preview: str = "") -> None:
        self.heartbeat += 1
        self.last_tool = tool_name
        self.last_args = args_preview[:80] if args_preview else ""



    def elapsed(self) -> str:
        secs = int(time.monotonic() - self.start_time)
        m, s = divmod(secs, 60)
        return f"{m}m{s:02d}s"


AGENT_PROGRESS = _AgentProgress()
_RECENT_FILE_ATTACHMENTS: dict[str, dict] = {}


def set_recent_file_attachments(attachments_by_alias: dict[str, dict]) -> None:
    """Register chat attachments that tools may reference by synthetic alias."""
    _RECENT_FILE_ATTACHMENTS.clear()
    _RECENT_FILE_ATTACHMENTS.update(attachments_by_alias or {})


def set_recent_image_attachments(attachments_by_alias: dict[str, dict]) -> None:
    """Backward-compatible wrapper for callers that only register images."""
    set_recent_file_attachments(attachments_by_alias)

# Set once at session start via set_project_path(); all tools use this as their workspace.
_PROJECT_PATH: str = ""
_PROJECT_SESSION = None
_PROJECT_STORE = None

# Toggled by run_skill() around a worker sub-agent's execution. Skill workers
# don't have the run_skill tool (only the orchestrator does), so read_file's
# large-file guidance must not tell a worker to "delegate via run_skill" —
# that advice is a dead end for it.
_IN_SKILL_WORKER: bool = False

# Whether the chat orchestrator currently holds the command-execution tools.
# Set from its model's policy when it is built: "direct" grants the workspace
# action tools (writes *and* execution), "delegate" withholds all of them. Kept
# next to _IN_SKILL_WORKER because both exist for the same reason -- recovery
# advice must name a route the caller can actually take (PROJECT_DESIGN 2.6) --
# but they answer different questions: one is "can I delegate?", the other is
# "can I run a command myself?".
_ORCHESTRATOR_HAS_TERMINAL: bool = True


def set_worker_context(active: bool) -> None:
    global _IN_SKILL_WORKER
    _IN_SKILL_WORKER = active


def set_orchestrator_terminal_access(active: bool) -> None:
    """Record whether the built orchestrator got the command-execution tools."""
    global _ORCHESTRATOR_HAS_TERMINAL
    _ORCHESTRATOR_HAS_TERMINAL = bool(active)


def caller_has_terminal() -> bool:
    """True when whoever is running right now can execute commands itself.

    A skill worker always can; the orchestrator only under the "direct" policy.
    Guidance branches on this instead of on the role, because after the direct
    policy started granting run_command the role no longer implies the answer.
    """
    return True if _IN_SKILL_WORKER else _ORCHESTRATOR_HAS_TERMINAL

def set_project_context(session, store=None) -> None:
    global _PROJECT_PATH, _PROJECT_SESSION, _PROJECT_STORE
    _PROJECT_SESSION = session
    _PROJECT_STORE = store
    if session:
        _PROJECT_PATH = os.path.abspath(session.project_path) if getattr(session, "project_path", "") else os.getcwd()
        # Snapshot the mode at context-set time so create_plan() can reference it
        # in audit messages even after the mode has already changed.
        session._initial_mode = getattr(session, "mode", "auto")

        # Ensure .opalatex/ is ignored by the user's own git so internal
        # files (editor state, skill requests) never appear in their git status.
        try:
            proj_gitignore = os.path.join(_PROJECT_PATH, ".gitignore")
            entry = ".opalatex/"
            existing = ""
            if os.path.isfile(proj_gitignore):
                with open(proj_gitignore, "r", encoding="utf-8") as _f:
                    existing = _f.read()
            if entry not in existing:
                with open(proj_gitignore, "a", encoding="utf-8") as _f:
                    if existing and not existing.endswith("\n"):
                        _f.write("\n")
                    _f.write(f"{entry}\n")
                # Untrack any .opalatex/ files already committed to the user's git index.
                subprocess.run(
                    ["git", "rm", "--cached", "-r", "--ignore-unmatch", ".opalatex/"],
                    cwd=_PROJECT_PATH, capture_output=True
                )
        except Exception:
            pass

        # Load project-specific .env file if it exists
        env_path = os.path.join(_PROJECT_PATH, ".env")
        if os.path.isfile(env_path):
            from dotenv import load_dotenv
            try:
                load_dotenv(dotenv_path=env_path, override=True)
            except Exception:
                pass
                
        # Also explicitly propagate api_key and api_base from session if present
        model_name = getattr(session, "model", None)
        alt_model_name = getattr(session, "worker_model", None)
        
        from .api_keys import get_env_var_for_model
        orch_env_var = get_env_var_for_model(model_name) if model_name else None
        worker_env_var = get_env_var_for_model(alt_model_name) if alt_model_name else None

        def env_file_has_var(var_name):
            if os.path.isfile(env_path):
                try:
                    with open(env_path, "r", encoding="utf-8") as f:
                        for line in f:
                            if line.strip().startswith(f"{var_name}="):
                                return True
                except Exception:
                    pass
            return False

        # 1. Orchestrator credentials
        if getattr(session, "api_key", None):
            os.environ["OPENAI_API_KEY"] = session.api_key
            if orch_env_var:
                os.environ[orch_env_var] = session.api_key
        else:
            if not env_file_has_var("OPENAI_API_KEY"):
                os.environ.pop("OPENAI_API_KEY", None)
            if orch_env_var and not env_file_has_var(orch_env_var):
                if not (orch_env_var == worker_env_var and getattr(session, "worker_api_key", None)):
                    os.environ.pop(orch_env_var, None)

        if getattr(session, "api_base", None):
            os.environ["OPENAI_API_BASE"] = session.api_base
        else:
            if not env_file_has_var("OPENAI_API_BASE"):
                os.environ.pop("OPENAI_API_BASE", None)

        # 2. Worker credentials
        if getattr(session, "worker_api_key", None):
            os.environ["WORKER_API_KEY"] = session.worker_api_key
            if worker_env_var:
                os.environ[worker_env_var] = session.worker_api_key
        else:
            if not env_file_has_var("WORKER_API_KEY"):
                os.environ.pop("WORKER_API_KEY", None)
            if worker_env_var:
                if not (worker_env_var == orch_env_var and getattr(session, "api_key", None)) and not env_file_has_var(worker_env_var):
                    os.environ.pop(worker_env_var, None)

        if getattr(session, "worker_api_base", None):
            os.environ["WORKER_API_BASE"] = session.worker_api_base
        else:
            if not env_file_has_var("WORKER_API_BASE"):
                os.environ.pop("WORKER_API_BASE", None)

        # Dynamically inject/register the configured context window (num_ctx) in LiteLLM's model mapping
        # so it doesn't fail with a BadRequestError (request exceeds available context window).
        try:
            import litellm
            from .config import resolve_effective_num_ctx
            model_name = getattr(session, "model", None)
            num_ctx = resolve_effective_num_ctx("memgpt", model_name, getattr(session, "api_base", None))
            if model_name and num_ctx:
                # Direct registration for model in litellm's dynamic model database
                if model_name not in litellm.model_prices_and_context_window:
                    litellm.model_prices_and_context_window[model_name] = {
                        "max_tokens": num_ctx,
                        "max_input_tokens": num_ctx,
                        "max_output_tokens": num_ctx,
                        "context_window": num_ctx,
                    }
                else:
                    litellm.model_prices_and_context_window[model_name]["context_window"] = num_ctx
                    litellm.model_prices_and_context_window[model_name]["max_tokens"] = num_ctx
                    litellm.model_prices_and_context_window[model_name]["max_input_tokens"] = num_ctx
        except Exception:
            pass
    else:
        _PROJECT_PATH = os.getcwd()


def get_project_path() -> str:
    return _PROJECT_PATH or os.getcwd()


def _collect_python_files(root: str, base_dir: str | None = None) -> list[str]:
    """Collect Python files below root while skipping OpalaTex internals."""
    base = os.path.abspath(base_dir or root)
    start = os.path.abspath(root)
    skipped_dirs = {"opalatex", "tests", "skills", "debug", "__pycache__", ".git", ".venv"}
    files: list[str] = []
    for dirpath, dirnames, filenames in os.walk(start):
        dirnames[:] = [d for d in dirnames if d not in skipped_dirs]
        rel_dir = os.path.relpath(dirpath, base)
        if rel_dir != "." and any(part in skipped_dirs for part in rel_dir.split(os.sep)):
            continue
        for filename in filenames:
            if filename.endswith(".py"):
                files.append(os.path.join(dirpath, filename))
    return files


def _resolve_path(path: str) -> str:
    """Make path absolute, rooted at the project directory if relative."""
    if os.path.isabs(path):
        return path
    return os.path.join(get_project_path(), path)


def _preview(value: object, max_len: int = 60) -> str:
    """Return a short, single-line preview of any argument value."""
    s = str(value).replace("\n", " ")
    return s[:max_len] + "…" if len(s) > max_len else s

_TEXT_FILE_ENCODINGS = ("utf-8-sig", "utf-8", "cp1252", "latin-1")


def _read_text_file(path: str) -> str:
    """Read text files with common encoding fallbacks used by LaTeX tools."""
    with open(path, "rb") as f:
        data = f.read()

    if not data:
        return ""

    encodings = list(_TEXT_FILE_ENCODINGS)
    if data.startswith(b"\xff\xfe"):
        encodings.insert(0, "utf-16-le")
    elif data.startswith(b"\xfe\xff"):
        encodings.insert(0, "utf-16-be")

    tried = set()
    for encoding in encodings:
        if encoding in tried:
            continue
        tried.add(encoding)
        try:
            text = data.decode(encoding)
            return text.replace("\r\n", "\n").replace("\r", "\n")
        except UnicodeDecodeError:
            continue

    text = data.decode("utf-8", errors="replace")
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _decode_escape_sequences(s: str) -> str:
    """Fix content where the model emitted literal \\n \\t \\r instead of real control chars.

    When a model double-escapes inside a JSON tool argument (\\\\n instead of \\n),
    json.loads produces the two-character sequence backslash-n. We detect this heuristic:
    the string has no real newlines but does contain literal backslash-n sequences,
    indicating the entire content is single-line with escaped line breaks.

    We do a simple targeted replacement of the common escape sequences to avoid
    corrupting multi-byte UTF-8 characters (accents, etc.) that unicode_escape would mangle.
    """
    if "\n" in s:
        return s
    if r"\n" not in s and r"\t" not in s and r"\r" not in s:
        return s
    return (
        s
        .replace(r"\n", "\n")
        .replace(r"\t", "\t")
        .replace(r"\r", "\r")
        .replace(r"\\", "\\")
    )

SAFE_SHELL_COMMANDS = {
    # Linux/macOS
    "ls", "pwd", "cat", "grep", "find", "whoami", "head", "tail", "less", "more", "tree", "cd", "echo",
    # Windows (CMD / PowerShell)
    "dir", "type", "findstr", "Get-ChildItem", "Get-Location", "Get-Content", "Select-String", "tasklist"
}

_DENIED_TOOLS = set()

class UserCancelException(BaseException):
    """Raised to immediately abort the agent turn when a user denies permission."""
    pass


def opalatex_tool(name: str, description: str, is_safe: bool = False):
    """
    Decorator that wraps the agenticblocks tool.
    Enforces 'plan', 'edit', and 'auto' mode safety rules.
    """
    def decorator(func):
        is_async = asyncio.iscoroutinefunction(func)

        def _check_permission(*args, **kwargs):
            mode = getattr(_PROJECT_SESSION, "mode", "auto") if _PROJECT_SESSION else "auto"
            
            # Auto mode or safe tool -> always execute
            if mode == "auto" or is_safe:
                return True
                
            # Check if this tool was set to 'always allow'
            if _PROJECT_SESSION and hasattr(_PROJECT_SESSION, "results"):
                allowed_tools = _PROJECT_SESSION.results.get("allowed_tools", [])
                if name in allowed_tools:
                    return True

            if name in _DENIED_TOOLS:
                return f"Execution blocked: You already asked to use '{name}' and the user DENIED it. Do not try again."

            # Now we are in 'edit' or 'plan' mode, and the tool is NOT safe.
            is_run_command_safe = False
            question = f"O agente quer usar a ferramenta '{name}'. Permitir?"
            
            if name == "run_command":
                cmd = kwargs.get("command") or (args[0] if args else "")
                base_cmd = cmd.strip().split()[0] if cmd else ""
                is_writing = ">" in cmd or "|" in cmd
                if base_cmd in SAFE_SHELL_COMMANDS and not is_writing:
                    is_run_command_safe = True
                else:
                    question = f"O agente quer executar o comando: '{cmd}'. Permitir?"
            
            if is_run_command_safe:
                return True
                
            if mode == "plan":
                return "Execution blocked: In 'plan' mode, you are not allowed to execute operations that modify the system."
                
            if mode == "edit":
                return question
                
            return True

        @functools.wraps(func)
        async def async_wrapper(*args, **kwargs):
            res = _check_permission(*args, **kwargs)
            if isinstance(res, str):
                if res.startswith("Execution blocked"):
                    return res
                
                # Ask user for permission using a custom confirm modal
                import uuid
                import asyncio
                from opalatex.agent_stdin import print_event, _gui_input_pending
                loop = asyncio.get_event_loop()
                req_id = str(uuid.uuid4())
                fut = loop.create_future()
                _gui_input_pending[req_id] = fut
                
                print_event("input_request", {
                    "id": req_id,
                    "prompt": res,
                    "type": "confirm",
                    "options": ["yes", "no", "always"],
                    "default": "yes"
                })
                
                try:
                    raw = await asyncio.wait_for(asyncio.shield(fut), timeout=86400.0) # wait up to 24h
                    ans = str(raw).strip().lower()
                except asyncio.TimeoutError:
                    ans = "no"
                finally:
                    _gui_input_pending.pop(req_id, None)
                    
                if ans not in ("yes", "y", "s", "sim", "true", "1", "always", "a"):
                    _DENIED_TOOLS.add(name)
                    return f"Execution blocked: The user denied permission to use '{name}'. Choose a safe alternative or explain what is needed."
                    
                if ans in ("always", "a"):
                    if _PROJECT_SESSION and hasattr(_PROJECT_SESSION, "results"):
                        if "allowed_tools" not in _PROJECT_SESSION.results:
                            _PROJECT_SESSION.results["allowed_tools"] = []
                        if name not in _PROJECT_SESSION.results["allowed_tools"]:
                            _PROJECT_SESSION.results["allowed_tools"].append(name)

            try:
                if is_async:
                    result = await func(*args, **kwargs)
                else:
                    result = func(*args, **kwargs)
            except Exception:
                raise

            return result
                
        return as_tool(name=name, description=description)(async_wrapper)
            
    return decorator

# ─── Tools ───────────────────────────────────────────────────────────────────
_BINARY_EXTS: set[str] = {
    ".xlsx", ".xls", ".ods", ".numbers",
    ".pdf", ".docx", ".pptx", ".odt", ".odp", ".doc", ".ppt", ".rtf",
    ".zip", ".rar", ".7z", ".gz", ".bz2", ".xz", ".tar", ".snap", ".whl", ".jar",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".webp", ".tiff", ".psd",
    ".mp3", ".mp4", ".wav", ".ogg", ".avi", ".mov", ".mkv", ".flac",
    ".exe", ".dll", ".so", ".dylib", ".bin", ".pyc", ".pyd", ".class",
    ".ttf", ".otf", ".woff", ".woff2", ".eot",
    ".db", ".sqlite", ".sqlite3", ".mdb",
}

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".webp", ".tiff", ".psd"}

# Binary formats whose text OpalaTex can extract by path (attachments.py). These
# never reach the binary refusal below: read_file extracts them and then treats
# the result like any other text, staged read included.
_DOC_EXTS = {".pdf", ".docx", ".pptx", ".xlsx"}

# Spreadsheet formats with no reader at all -- .xlsx is handled above.
_SHEET_EXTS = {".xls", ".ods", ".numbers"}


def _binary_read_error(resolved: str) -> str | None:
    """Return a diagnostic when *resolved* is not text, else None.

    `_read_text_file` falls back to latin-1, which decodes *any* byte sequence
    without raising — so without this check a .xlsx or .pdf read by path returns
    mojibake that looks like content and lands in the context window. Failing
    with a route the caller can actually take is the contract (PROJECT_DESIGN
    2.6, no silent substitution).
    """
    ext = os.path.splitext(resolved)[1].lower()

    binary = ext in _BINARY_EXTS
    if not binary:
        try:
            with open(resolved, "rb") as f:
                head = f.read(8192)
        except OSError:
            return None
        # UTF-16/32 text is full of NUL bytes and _read_text_file handles it via
        # the BOM, so a BOM rules out the NUL sniff.
        if not head.startswith((b"\xff\xfe", b"\xfe\xff")) and b"\x00" in head:
            binary = True

    if not binary:
        return None

    preview = _preview(resolved)

    if ext in _IMAGE_EXTS:
        route = f"Use analyze_image('{preview}', prompt) to inspect it."
    else:
        if ext in _SHEET_EXTS:
            what = f"OpalaTex reads .xlsx but has no reader for {ext}."
            conversion = "an .xlsx or .csv export"
        else:
            what = "This is not a text file."
            conversion = "a conversion to text"

        # Naming a route the caller cannot take is a dead end that pushes models
        # into improvising (PROJECT_DESIGN 2.6). A worker always has run_command
        # and so does a "direct" orchestrator; only a "delegate" orchestrator has
        # to route this through a worker, which plan mode blocks outright --
        # leaving asking the user as its only move.
        if caller_has_terminal():
            how = (
                f"You have run_command: convert it yourself with {conversion} into a text file "
                "inside the project, then read that file."
            )
        else:
            how = (
                f"You have no terminal tool, so you cannot convert it yourself. Delegate {conversion} "
                "to the 'command-line' skill with run_skill and then read the converted file. If you "
                "are in plan mode, run_skill is blocked: ask the user to attach the file to the chat "
                "or to provide the converted text instead of planning around content you cannot read."
            )
        route = f"{what} {how} Do not retry read_file on this path."

    return f"Error: '{preview}' is a binary file ({ext or 'no extension'}), not text. {route}"


@opalatex_tool(
    name="read_file",
    is_safe=True,
    description=(
        "Read the contents of a file in the project workspace. Relative paths are resolved "
        "from the project directory. PDF, DOCX, PPTX and XLSX files are supported: their text "
        "is extracted automatically, so read them with this tool instead of asking the user to "
        "convert them by hand."
    ),
)
def read_file(path: str) -> str:
    attachment = _RECENT_FILE_ATTACHMENTS.get(path) or _RECENT_FILE_ATTACHMENTS.get(os.path.basename(path))
    if attachment:
        att_type = attachment.get("type", "")
        if att_type == "pdf_text":
            # Uploads made from this version cache the original on disk and carry
            # a path; history written before that still carries the base64 inline.
            raw_path = attachment.get("raw_path")
            raw_data = attachment.get("raw_data")
            if raw_path and os.path.exists(raw_path):
                try:
                    from opalatex.attachments import extract_document_text_from_path
                    return extract_document_text_from_path(raw_path)
                except Exception as e:
                    raise ValueError(f"Error extracting attached document '{attachment.get('name', path)}': {e}")
            if raw_data:
                try:
                    from opalatex.attachments import extract_document_text
                    return extract_document_text(
                        raw_data,
                        attachment.get("raw_mime") or attachment.get("mime", ""),
                        attachment.get("name", path),
                    )
                except Exception as e:
                    raise ValueError(f"Error extracting attached document '{attachment.get('name', path)}': {e}")
            # No original available (old chat whose cache entry is gone): the
            # extracted text is exactly what the model was shown for this turn.
            return attachment.get("data", "")
        if att_type == "image":
            raise ValueError(
                f"'{path}' is an image attachment. Use analyze_image('{path}', prompt) to inspect it."
            )

    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update("read_file", f"path={_preview(resolved)}")

    try:
        if os.path.isdir(resolved):
            raise ValueError(f"Error: '{_preview(resolved)}' is a directory, not a file. Use run_command with 'ls' or get_project_overview() to view contents.")
        if not os.path.exists(resolved):
            raise ValueError(f"Error: file not found: {_preview(resolved)}. If you are trying to write code, use 'write_file' instead.")
    except OSError as e:
        # Catch cases like [Errno 36] File name too long if path contains code
        raise ValueError(f"Error: invalid path argument ({e.strerror}). 'read_file' expects a file path, not file contents.")

    # A PDF/DOCX/PPTX/XLSX in the workspace is not text, but it does have text
    # inside it. Extract it here so the rest of read_file (the context budget
    # check) treats it exactly like any other document, instead of refusing a
    # file whose content OpalaTex can actually produce.
    extracted = None
    if os.path.splitext(resolved)[1].lower() in _DOC_EXTS:
        from .attachments import extract_document_text_from_path

        AGENT_PROGRESS.update("read_file", f"extracting {_preview(resolved)}")
        try:
            extracted = extract_document_text_from_path(resolved)
        except ImportError as e:
            raise ValueError(
                f"Error: '{_preview(resolved)}' needs a document extractor that is not "
                f"installed ({e}). Ask the user to install it, or to provide a text export."
            )
        except Exception as e:
            raise ValueError(
                f"Error extracting text from '{_preview(resolved)}': {e}. The file may be "
                "corrupt, password-protected, or empty. Do not retry read_file on this path."
            )
        if not extracted.strip():
            raise ValueError(
                f"Error: '{_preview(resolved)}' was read but contains no extractable text "
                "(a scanned/image-only PDF, or an empty document). Do not retry read_file on "
                "this path; use analyze_image on the page images, or ask the user for a text version."
            )

    if extracted is None:
        binary_error = _binary_read_error(resolved)
        if binary_error:
            raise ValueError(binary_error)
    else:
        budget_chars = free_context_chars()
        if budget_chars > 0 and len(extracted) > budget_chars:
            # read_content_pos cannot page this: it would read the raw bytes of a
            # binary container. The reachable route is to land the extracted text
            # in the project as a text file and work on that.
            if caller_has_terminal():
                how = (
                    "Use run_command to extract it to a .txt/.md file in the project "
                    "(pymupdf4llm for PDF, openpyxl for XLSX, python-docx/python-pptx are "
                    "installed), then search_code + read_content_pos that file."
                )
            else:
                how = (
                    "Delegate to the 'command-line' skill with run_skill: have it extract the "
                    "document to a .txt/.md file in the project, then use search_code + "
                    "read_content_pos on that file. In plan mode run_skill is blocked, so ask "
                    "the user for the part they need instead."
                )
            raise ValueError(
                f"Error: the text extracted from '{_preview(resolved)}' is {len(extracted):,} "
                f"characters and does not fit the remaining context budget "
                f"(~{budget_chars:,} characters). Do not retry read_file on this path, and do "
                f"not call read_content_pos on it — it is a binary container. {how}"
            )
        header = (
            f"[Text extracted from '{os.path.basename(resolved)}' "
            f"({os.path.splitext(resolved)[1].lower().lstrip('.')}). Layout and formatting are "
            "not preserved.]\n"
        )
        return header + extracted

    # A whole-file read is the fastest way to destroy a context window: the
    # result is appended to history, MemGPT cannot evict the current user turn
    # to make room, and the provider then truncates the request from the front —
    # silently dropping the very question being answered. Refuse the read with a
    # diagnostic instead of returning content that cannot fit (PROJECT_DESIGN
    # 2.6, no silent substitution): the model is told how to page through the
    # file rather than having its request quietly turned into a partial one.
    budget_chars = free_context_chars()
    try:
        size = os.path.getsize(resolved)
    except OSError:
        size = 0
    if size > budget_chars:
        window = context_window_tokens()
        used = used_context_tokens()
        if budget_chars <= 0:
            raise ValueError(
                f"Error: the context window is exhausted ({used:,} of {window:,} tokens "
                f"already used), so '{_preview(resolved)}' cannot be read. Summarize what "
                "you already know and answer, or ask the user to start a new chat or raise "
                "num_ctx in the project settings."
            )
        # Structured data/log files should be routed to a specialized skill
        # rather than paged through manually.  The redirect is strongest here
        # because weak models obey error-message guidance more reliably than
        # system-prompt instructions.
        _DATA_EXTS = {".jsonl", ".csv", ".tsv", ".log"}
        ext = os.path.splitext(resolved)[1].lower()
        if ext in _DATA_EXTS:
            if _IN_SKILL_WORKER:
                # You are the worker, not the orchestrator: you have no run_skill
                # tool, so telling you to "delegate via run_skill" is a dead end.
                routing_hint = (
                    "IMPORTANT: You are a skill worker and do NOT have a run_skill tool. "
                    "Use the processing script or command listed in your skill instructions "
                    "(via run_command) to inspect or condense this file instead of reading it "
                    "directly. If your skill has no such script, report back that the file is "
                    "too large to process directly."
                )
            else:
                routing_hint = (
                    "IMPORTANT: Check your Available skills list for a specialized "
                    "log/data processing skill and delegate to it with run_skill. "
                    "If no such skill is active, use the command-line skill to run "
                    "a small streaming Python script that processes the file."
                )
            raise ValueError(
                f"Error: '{_preview(resolved)}' is a large data/log file "
                f"({size:,} bytes, ~{size // 4:,} tokens) and does not fit "
                f"the remaining context budget (~{budget_chars // 4:,} tokens "
                f"of the {window:,}-token window, {used:,} already used). "
                f"Do NOT retry read_file or read_content_pos on this path. "
                f"{routing_hint}"
            )
        raise ValueError(
            f"Error: '{_preview(resolved)}' is {size:,} bytes (~{size // 4:,} tokens) and does "
            f"not fit the remaining context budget (~{budget_chars // 4:,} tokens of the "
            f"{window:,}-token window, {used:,} already used). Do not retry read_file on this "
            f"path. Use search_code(pattern, path) to locate the relevant lines, then "
            f"read_content_pos(path, start_pos, end_pos) to read only that range."
        )

    try:
        return _read_text_file(resolved)
    except Exception as e:
        raise ValueError(f"Error reading {_preview(resolved)}: {e}")

@opalatex_tool(name="analyze_image", is_safe=True, description="Analyzes an image file using computer vision and returns a text description. Use this when the user asks you to look at a chart, UI mockup, or any image file.")
def analyze_image(file_path: str, prompt: str = "Describe this image in detail") -> str:
    model = ""
    try:
        import base64
        import mimetypes
        import litellm
        from .config import get_agent_llm_kwargs, get_agent_model, sanitize_litellm_kwargs_for_model
        
        attachment = _RECENT_FILE_ATTACHMENTS.get(file_path) or _RECENT_FILE_ATTACHMENTS.get(os.path.basename(file_path))
        if attachment:
            resolved = file_path
            encoded_img = attachment.get("data", "")
            content_type = attachment.get("mime") or "image/jpeg"
        else:
            resolved = _resolve_path(file_path)
        if not attachment and (not os.path.exists(resolved) or os.path.isdir(resolved)):
            return "infelizmente não consegui analisar sua imagem: arquivo não encontrado ou é um diretório."
            
        if not attachment:
            content_type, _ = mimetypes.guess_type(resolved)
            if not content_type or not content_type.startswith("image/"):
                content_type = "image/jpeg"
            
        if not attachment:
            with open(resolved, "rb") as f:
                encoded_img = base64.b64encode(f.read()).decode("utf-8")
            
        kwargs = get_agent_llm_kwargs("memgpt")
        
        # Use the orchestrator's model by default (since it's the one using the tool)
        model = get_agent_model("memgpt")
        if _PROJECT_SESSION:
            project_override = getattr(_PROJECT_SESSION, "model", None)
            if project_override:
                model = project_override
                
        kwargs.pop("stream", None)
        
        if kwargs.get("api_base"):
            if model.startswith("ollama/") or model.startswith("ollama_chat/"):
                if kwargs["api_base"].endswith("/v1"):
                    kwargs["api_base"] = kwargs["api_base"][:-3]
                elif kwargs["api_base"].endswith("/v1/"):
                    kwargs["api_base"] = kwargs["api_base"][:-4]
        kwargs = sanitize_litellm_kwargs_for_model(model, kwargs)
        kwargs.pop("stream", None)

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{content_type};base64,{encoded_img}"}
                    }
                ]
            }
        ]
        
        AGENT_PROGRESS.update("analyze_image", f"path={_preview(resolved)}")
        response = litellm.completion(model=model, messages=messages, **kwargs)
        return response.choices[0].message.content or "infelizmente não consegui analisar sua imagem."
        
    except litellm.exceptions.BadRequestError as e:
        return f"CRITICAL ERROR: The model ({model}) rejected the request. Explicitly inform the user that this model likely does not support Computer Vision (not multimodal) or the image is too large. Technical detail: {str(e)}"
    except litellm.exceptions.AuthenticationError as e:
        return f"CRITICAL ERROR: Authentication failed for model {model}. EXPLÍCITLY tell the user that they need to specify the API Key for this model in the top bar (under 'Edit Models'). Technical detail: {str(e)}"
    except litellm.exceptions.APIConnectionError as e:
        return f"CRITICAL ERROR: Connection failed with the model's API ({model}). Tell the user to verify if the local server (e.g., Ollama) is running or if the URL is correct. Technical detail: {str(e)}"
    except Exception as e:
        return f"CRITICAL ERROR: An internal error occurred while analyzing the image using the model {model}. Tell the user the following reason: {str(e)}"


# ─── Image generation ─────────────────────────────────────────────────────────
# One tool, several providers: the transport lives in
# `agenticblocks.blocks.image`, and which provider answers is a catalog entry
# plus an api_base (Settings > Image Generation). The tool itself only knows
# "prompt in, file on disk out".

_IMAGE_ERROR_GUIDANCE = {
    "auth": (
        "Tell the user the image model's API key is missing or was rejected, and "
        "that they can set it in the top bar under 'Edit Models'."
    ),
    "connection": (
        "Tell the user the image endpoint could not be reached: if it is a local "
        "server (LocalAI, Docker Model Runner, vLLM-Omni, Ollama), it must be "
        "running and its api_base must match what is registered in 'Edit Models'."
    ),
    "unsupported": (
        "Tell the user this endpoint has no image generation route. A local "
        "Ollama only serves one on builds where image generation is available "
        "(currently macOS); on Linux/Windows use LocalAI, Docker Model Runner or "
        "vLLM-Omni, or point api_base at a machine that does serve it."
    ),
    "bad_request": (
        "The provider rejected the request. Report the technical detail to the "
        "user: usually the model id or the requested size is not supported."
    ),
    "empty": (
        "The provider accepted the request but returned no image. This is often a "
        "content filter. Report it to the user instead of retrying the same prompt."
    ),
    "unknown_route": (
        "The configured image route does not exist. Tell the user to review "
        "Settings > Image Generation."
    ),
    "not_configured": (
        "Tell the user to configure an image model in Settings > Image Generation."
    ),
}


def _image_filename_stem(filename: str, prompt: str) -> str:
    """Return a safe file stem from *filename*, or one derived from *prompt*."""
    raw = (filename or "").strip().replace("\\", "/")
    raw = raw.split("/")[-1]
    raw = os.path.splitext(raw)[0]
    if not raw:
        raw = "-".join((prompt or "image").lower().split())[:48]
    stem = re.sub(r"[^A-Za-z0-9._-]+", "-", raw).strip("-._")
    return stem or "image"


def _unique_image_path(directory: str, stem: str, ext: str) -> str:
    """Return a path inside *directory* that does not overwrite an existing file."""
    candidate = os.path.join(directory, f"{stem}{ext}")
    index = 2
    while os.path.exists(candidate):
        candidate = os.path.join(directory, f"{stem}-{index}{ext}")
        index += 1
    return candidate


def _format_byte_size(count: int) -> str:
    """Return a human-readable size; small files must not read as "0 KB"."""
    if count < 1024:
        return f"{count} bytes"
    if count < 1024 * 1024:
        return f"{count / 1024:.0f} KB"
    return f"{count / (1024 * 1024):.1f} MB"


def _write_generated_image(target_path: str, data: bytes) -> tuple[int, str]:
    """Write *data* to *target_path*; return (byte count, "WxH" or "")."""
    Path(target_path).parent.mkdir(parents=True, exist_ok=True)
    with open(target_path, "wb") as f:
        f.write(data)

    dimensions = ""
    try:
        from PIL import Image

        with Image.open(target_path) as img:
            dimensions = f"{img.width}x{img.height}"
    except Exception:
        pass
    return len(data), dimensions


@opalatex_tool(
    name="generate_image",
    is_safe=False,
    description=(
        "Generate an image from a text prompt and save it as a file inside the "
        "project, returning its project-relative path. Use it when the user asks "
        "for an illustration, figure or diagram that does not exist yet. The "
        "prompt should be a detailed visual description in English. Never use it "
        "to render LaTeX, plots or charts -- those belong in TikZ/pgfplots code. "
        "The returned path can be shown in chat as Markdown (![caption](path)) "
        "and included in LaTeX with \\includegraphics."
    ),
)
async def generate_image(
    prompt: str,
    filename: str = "",
    size: str = "",
    negative_prompt: str = "",
    seed: int = -1,
) -> str:
    from agenticblocks.blocks.image import (
        ImageGenerationError,
        ImageGenerationInput,
        extension_for_mime,
        resolve_image_bytes,
    )

    from .image_gen_config import build_block, configuration_problem, load_config

    if not (prompt or "").strip():
        return "CRITICAL ERROR: generate_image requires a non-empty 'prompt'."

    problem = configuration_problem()
    if problem:
        return f"CRITICAL ERROR: {problem}"

    cfg = load_config()
    AGENT_PROGRESS.update("generate_image", f"prompt={_preview(prompt)}")

    try:
        block = build_block()
        request = ImageGenerationInput(
            prompt=prompt,
            n=1,
            size=(size or cfg.get("size") or "").strip(),
            negative_prompt=(negative_prompt or "").strip(),
            seed=(int(seed) if seed is not None and int(seed) >= 0 else None),
        )
        output = await block.run(request)
        data, mime = await resolve_image_bytes(output.images[0])
    except ImageGenerationError as e:
        guidance = _IMAGE_ERROR_GUIDANCE.get(e.kind, "Report the failure to the user.")
        return f"CRITICAL ERROR: {e} {guidance}"
    except Exception as e:
        return (
            "CRITICAL ERROR: An internal error occurred while generating the image. "
            f"Tell the user the following reason: {e}"
        )

    output_dir = str(cfg.get("output_dir") or "figures")
    target_dir = _resolve_path(output_dir)
    stem = _image_filename_stem(filename, prompt)
    target_path = _unique_image_path(target_dir, stem, extension_for_mime(mime))

    try:
        byte_count, dimensions = await asyncio.to_thread(
            _write_generated_image, target_path, data
        )
    except OSError as e:
        raise ValueError(f"Error saving the generated image: {e}")

    relative = os.path.relpath(target_path, get_project_path()).replace(os.sep, "/")
    latex_path = os.path.splitext(relative)[0]
    detail = f"{dimensions}, " if dimensions else ""

    # The bytes never enter the conversation: a 1024x1024 PNG in base64 is over a
    # megabyte of context, and the file on disk is the deliverable anyway
    # (PROJECT_DESIGN 2.6, tool results are bounded by the remaining window).
    return (
        f"Image generated with {output.model} and saved to {relative} "
        f"({detail}{_format_byte_size(byte_count)}).\n"
        f"To show it to the user in chat, write the Markdown ![{stem}]({relative}).\n"
        f"To place it in a LaTeX document, use \\includegraphics[width=0.8\\textwidth]{{{latex_path}}} "
        f"inside a figure environment."
    )

@opalatex_tool(name="write_file", is_safe=False, description="Write or overwrite a file inside the project directory. Relative paths are resolved from the project directory. Creates parent directories if needed. ALWAYS use this tool to save file content — never use run_command with echo/printf/cat to write files, as shell quoting will break with multi-line or HTML/JSON content.")
def write_file(path: str, content: str) -> str:
    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update("write_file", f"path={_preview(resolved)}")
    
    try:
        if os.path.isdir(resolved):
            raise ValueError(f"Error: '{_preview(resolved)}' is a directory. Cannot write to a directory.")
        Path(resolved).parent.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        raise ValueError(f"Error: invalid path argument ({e.strerror}).")

    try:
        content = _decode_escape_sequences(content)
        with open(resolved, "w", encoding="utf-8", newline="") as f:
            f.write(content)
        try:
            from .code_index import CODE_INDEX
            CODE_INDEX.rebuild_file(resolved)
        except Exception:
            pass
            
        return f"Successfully wrote to {_preview(resolved)}."
    except Exception as e:
        raise ValueError(f"Error writing {_preview(resolved)}: {e}")


@opalatex_tool(
    name="export_tex_to_docx",
    is_safe=False,
    description=(
        "Export a .tex file in the project to .docx using Pandoc. "
        "Use this when the user asks to convert LaTeX/TeX source to Word. "
        "The input and output paths are project-relative. Pandoc must be installed."
    ),
)
def export_tex_to_docx(tex_path: str, output_path: str = "") -> str:
    resolved_project = get_project_path()
    AGENT_PROGRESS.update("export_tex_to_docx", f"path={_preview(tex_path)}")
    try:
        from .document_exporter import export_tex_to_docx as _export_tex_to_docx
        result = _export_tex_to_docx(resolved_project, tex_path, output_path)
    except Exception as e:
        raise ValueError(f"Error exporting DOCX: {e}")
    if not result.get("success"):
        raise ValueError(result.get("log") or "DOCX export failed.")
    return f"Successfully exported DOCX to {result.get('relative_output_path') or result.get('output_path')}."


def _require_extension(path: str, allowed: set[str]) -> None:
    ext = Path(path).suffix.lower()
    if ext not in allowed:
        raise ValueError(f"Output path must end with one of: {', '.join(sorted(allowed))}")


@opalatex_tool(
    name="create_docx_file",
    is_safe=False,
    description=(
        "Create a Word .docx file directly from simple Markdown-like text. "
        "Supports headings beginning with #, ##, ###, bullet lines beginning with '- ', "
        "numbered lines like '1. ', and plain paragraphs. Use this instead of writing raw binary DOCX."
    ),
)
def create_docx_file(path: str, markdown: str, title: str = "") -> str:
    try:
        resolved = _resolve_path(path)
        _require_extension(resolved, {".docx"})
    except Exception as e:
        raise ValueError(f"Error resolving DOCX path: {e}")

    AGENT_PROGRESS.update("create_docx_file", f"path={_preview(resolved)}")
    try:
        from docx import Document
    except ImportError:
        raise ValueError("python-docx is not installed. Install project dependencies before creating DOCX files directly.")

    try:
        Path(resolved).parent.mkdir(parents=True, exist_ok=True)
        doc = Document()
        if title:
            doc.add_heading(title, level=1)
        for raw_line in _decode_escape_sequences(markdown).splitlines():
            line = raw_line.strip()
            if not line:
                continue
            if line.startswith("### "):
                doc.add_heading(line[4:].strip(), level=3)
            elif line.startswith("## "):
                doc.add_heading(line[3:].strip(), level=2)
            elif line.startswith("# "):
                doc.add_heading(line[2:].strip(), level=1)
            elif line.startswith("- "):
                doc.add_paragraph(line[2:].strip(), style="List Bullet")
            elif len(line) > 3 and line[0].isdigit() and line[1:3] == ". ":
                doc.add_paragraph(line[3:].strip(), style="List Number")
            else:
                doc.add_paragraph(line)
        doc.save(resolved)
        return f"Successfully created DOCX file at {_preview(resolved)}."
    except Exception as e:
        raise ValueError(f"Error creating DOCX file {_preview(resolved)}: {e}")


@opalatex_tool(
    name="create_pptx_file",
    is_safe=False,
    description=(
        "Create a PowerPoint .pptx file directly from a JSON slide outline. "
        "slides_json must be an array of objects like "
        "[{\"title\":\"Slide title\",\"bullets\":[\"Point one\",\"Point two\"]}]."
    ),
)
def create_pptx_file(path: str, slides_json: str, title: str = "") -> str:
    try:
        resolved = _resolve_path(path)
        _require_extension(resolved, {".pptx"})
    except Exception as e:
        raise ValueError(f"Error resolving PPTX path: {e}")

    AGENT_PROGRESS.update("create_pptx_file", f"path={_preview(resolved)}")
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("python-pptx is not installed. Install project dependencies before creating PPTX files directly.")

    try:
        slides = json.loads(_decode_escape_sequences(slides_json))
        if not isinstance(slides, list):
            raise ValueError("slides_json must be a JSON array")
        Path(resolved).parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        if title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = ""
        for item in slides:
            if not isinstance(item, dict):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(item.get("title", "Untitled slide"))
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = item.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [line.strip() for line in bullets.splitlines() if line.strip()]
            if not bullets:
                bullets = [str(item.get("body", ""))] if item.get("body") else []
            for index, bullet in enumerate(bullets):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
        prs.save(resolved)
        return f"Successfully created PPTX file at {_preview(resolved)}."
    except Exception as e:
        raise ValueError(f"Error creating PPTX file {_preview(resolved)}: {e}")


import platform
from .subprocess_utils import utf8_text_kwargs

_os_info = f"{platform.system()} ({platform.release()})"
_run_cmd_desc = f"Execute a non-interactive shell command (e.g. ls, dir, mkdir, npm install). Runs inside the project directory. Returns stdout/stderr. Use search_code for project text searches. NEVER run servers or infinite processes. NEVER use echo/printf/cat to write file content — use write_file instead. NOTE: Host OS is {_os_info}. Use the appropriate shell commands."

@opalatex_tool(name="run_command", is_safe=False, description=_run_cmd_desc)
def run_command(command: str) -> str:
    AGENT_PROGRESS.update("run_command", f"$ {_preview(command)}")
    cwd = get_project_path()
    try:
        res = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            **utf8_text_kwargs(),
            stdin=subprocess.DEVNULL,
            timeout=120,
            cwd=cwd,
        )
        out = res.stdout.strip()
        err = res.stderr.strip()
        # Truncate large outputs
        if len(out) > 2000:
            out = out[:1000] + "\n... [TRUNCATED] ...\n" + out[-500:]
        if len(err) > 2000:
            err = err[:1000] + "\n... [TRUNCATED] ...\n" + err[-500:]
        import json
        result = {
            "stdout": out,
            "stderr": err,
            "exit_code": res.returncode
        }
        return json.dumps({"result": json.dumps(result, indent=2)})
    except subprocess.TimeoutExpired:
        raise ValueError("Error: Command timed out after 120 seconds.")
    except Exception as e:
        raise ValueError(f"Error running command: {e}")

@opalatex_tool(name="run_background_command", is_safe=False, description="Start a long-running background command or server (e.g., `npm run dev`) directly in the user's MAIN IDE terminal. This runs asynchronously and returns immediately, allowing the agent to continue working while the server runs. NEVER use this for commands where you need to see the output to proceed.")
async def run_background_command(command: str) -> str:
    import json
    import urllib.request
    
    AGENT_PROGRESS.update("background_cmd", f"$ {_preview(command)}")
    cwd = get_project_path()

    try:
        req = urllib.request.Request(
            "http://127.0.0.1:3000/api/terminal/input",
            data=json.dumps({
                "action": "input",
                "text": f"{command}\r",
                "term_id": "main",
                "projectPath": cwd
            }).encode('utf-8'),
            headers={'Content-Type': 'application/json'},
            method='POST'
        )
        with urllib.request.urlopen(req, timeout=2.0) as response:
            if response.status == 200:
                return f"SUCCESS: The command '{command}' has been sent to the main background terminal and is now running."
            else:
                return f"FAILED to start background command: HTTP {response.status}"
    except Exception as e:
        return f"FAILED to send command to background terminal: {str(e)}"


@opalatex_tool(name="run_python_script", is_safe=False, description="Execute a Python script securely. It automatically uses the correct Python interpreter for the environment. Provide the script path and any optional arguments.")
def run_python_script(script_path: str, args: str = "") -> str:
    try:
        resolved = _resolve_path(script_path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    cmd = f'"{sys.executable}" "{resolved}" {args}'.strip()
    AGENT_PROGRESS.update("run_python_script", f"$ {_preview(cmd)}")
    cwd = get_project_path()
    try:
        res = subprocess.run(
            cmd,
            shell=True,
            capture_output=True,
            **utf8_text_kwargs(),
            stdin=subprocess.DEVNULL,
            timeout=120,
            cwd=cwd,
        )
        out = res.stdout.strip()
        err = res.stderr.strip()
        # Truncate large outputs
        if len(out) > 2000:
            out = out[:1000] + "\n... [TRUNCATED] ...\n" + out[-500:]
        if len(err) > 2000:
            err = err[:1000] + "\n... [TRUNCATED] ...\n" + err[-500:]
        import json
        result = {
            "stdout": out,
            "stderr": err,
            "exit_code": res.returncode
        }
        return json.dumps({"result": json.dumps(result, indent=2)})
    except subprocess.TimeoutExpired:
        raise ValueError("Error: Command timed out after 120 seconds.")
    except Exception as e:
        raise ValueError(f"Error running script: {e}")

@opalatex_tool(name="run_interactive_command", is_safe=False, description="Run a command that requires user interaction (e.g. npm create, interactive scripts, prompts). This opens a dedicated interactive terminal popup in the user's GUI, allowing them to answer the prompts safely. Use this WHENEVER a command needs human choices, waits for input, or requires a PTY. Do NOT use standard `exec` or `run_command` for interactive tasks, as they will hang.")
async def run_interactive_command(command: str) -> str:
    import uuid
    import opalatex.terminal as T
    
    AGENT_PROGRESS.update("interactive_cmd", _preview(command))
    term_id = "temp_" + str(uuid.uuid4())[:8]

    # Trigger the interactive terminal modal in the UI
    res = await T.a_interactive_terminal(command, term_id)
    
    if res and res.lower() not in ("cancel", "no", "false"):
        return f"SUCCESS: The interactive command '{command}' finished successfully according to the user."
    else:
        return f"FAILURE: The user indicated that the interactive command '{command}' failed or was cancelled. Please ask the user for details if needed."

@opalatex_tool(
    name="search_code",
    is_safe=True,
    description=(
        "Search for text or a regular expression in project files using Python, "
        "returning relative file paths, 1-indexed line numbers, and matching lines. "
        "Use this to locate sections, labels, definitions, or other markers before "
        "calling read_content_pos or replace_content_range."
    ),
)
def search_code(
    query: str,
    path: str = ".",
    use_regex: bool = False,
    case_sensitive: bool = False,
    max_results: int = 50,
    file_pattern: str = "*",
) -> str:
    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update("search_code", f"query={_preview(query)} path={_preview(resolved)}")

    if not query:
        raise ValueError("query must not be empty.")
    if max_results < 1:
        raise ValueError("max_results must be a positive integer.")

    root = Path(get_project_path()).resolve()
    target = Path(resolved).resolve()
    skipped_dirs = {".git", ".opalatex", "node_modules", "__pycache__", ".venv", "venv", "dist", "build"}
    skipped_exts = {
        ".aux", ".bbl", ".blg", ".class", ".dll", ".exe", ".gz", ".ico", ".jar",
        ".jpg", ".jpeg", ".log", ".mp3", ".mp4", ".pdf", ".png", ".pyc", ".so",
        ".synctex.gz", ".zip",
    }

    try:
        if target.is_file():
            candidates = [target]
        elif target.is_dir():
            candidates = []
            for dirpath, dirnames, filenames in os.walk(target):
                dirnames[:] = [d for d in dirnames if d not in skipped_dirs and not d.startswith(".")]
                for filename in filenames:
                    candidate = Path(dirpath) / filename
                    if candidate.suffix.lower() in skipped_exts:
                        continue
                    rel_candidate = os.path.relpath(candidate, target).replace(os.sep, "/")
                    if not (
                        fnmatch(filename, file_pattern)
                        or fnmatch(rel_candidate, file_pattern)
                    ):
                        continue
                    candidates.append(candidate)
        else:
            raise ValueError(f"Error: path does not exist: {_preview(resolved)}.")

        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = re.compile(query if use_regex else re.escape(query), flags)
        matches: list[str] = []
        searched = 0

        for candidate in candidates:
            searched += 1
            try:
                text = _read_text_file(str(candidate))
            except Exception:
                continue
            if "\x00" in text[:4096]:
                continue
            rel_path = os.path.relpath(candidate, root)
            for line_number, line in enumerate(text.splitlines(), start=1):
                if pattern.search(line):
                    matches.append(f"{rel_path}:{line_number}: {line.strip()}")
                    if len(matches) >= max_results:
                        return (
                            "\n".join(matches)
                            + f"\n[search_code] Stopped after {max_results} matches."
                        )

        if not matches:
            return f"No matches for {query!r} in {_preview(os.path.relpath(target, root))}."
        return "\n".join(matches) + f"\n[search_code] Searched {searched} files."
    except Exception as e:
        raise ValueError(f"Error searching code in {_preview(resolved)}: {e}")

@opalatex_tool(
    name="ask_question",
    is_safe=True,
    description=(
        "Ask the user a clarifying question with optional selectable choices. "
        "Provide 'question' and optionally a list of 2-5 distinct 'options' representing common choices. "
        "The user will be presented with clickable choices in the UI plus an automatic 'Other' write-in field. "
        "Do NOT include your own 'Other', 'Others (please specify)', or 'None of the above' choice in 'options' — "
        "that catch-all is added automatically by the UI, so adding one yourself creates a duplicate. "
        "Use this proactively whenever collecting user choices, column selections, target formats, seed filters, or clarifying ambiguous preferences."
    )
)
async def ask_question(
    question: str,
    options: list[str] | None = None,
    is_multi_select: bool = False
) -> str:
    AGENT_PROGRESS.update("ask_question", _preview(question))
    
    import opalatex.terminal as T
    ans = await T.aask(question, options=options, is_multi_select=is_multi_select)
    
    if not ans:
        return "The user did not provide an answer or cancelled the prompt."
    return f"User response: {ans}"

# Backward compatibility alias
ask_human = ask_question



@opalatex_tool(
    name="get_project_overview",
    is_safe=True,
    description=(
        "Return a compact overview of the current project: directory tree (max depth defined by parameter max_depth, DEFAULT/MINIMUM 5), "
        "file count by type, and a summary of key files (README, package.json, requirements.txt, etc.). "
        "Use this when the project structure or target file is unknown. Skip it when the user or context already identifies the file to inspect or edit."
    ),
)
def get_project_overview(max_depth:int = 5) -> str:
    from .config import get_project_overview_max_depth
    AGENT_PROGRESS.update("get_project_overview")
    root = Path(get_project_path())
    
    try:
        if max_depth < 0:
            max_depth_cfg = get_project_overview_max_depth()
        else:
            max_depth_cfg = max_depth
            
        # FORCE MINIMUM DEPTH
        if max_depth_cfg < 5:
            max_depth_cfg = 5

        if not root.exists():
            return f"Project directory does not exist yet: {root}"

        # ── Directory tree (skip hidden & node_modules/__pycache__) ──
        SKIP = {".git", "node_modules", "__pycache__", ".venv", "venv", ".mypy_cache", "dist", "build"}

        def _tree(path: Path, depth: int = 0) -> list[str]:
            if depth >= max_depth_cfg:
                return ["  " * depth + "..."]
            lines = []
            try:
                entries = sorted(path.iterdir(), key=lambda p: (p.is_file(), p.name))
            except PermissionError:
                return []
            for entry in entries:
                if entry.name.startswith(".") or entry.name in SKIP:
                    continue
                indent = "  " * depth
                lines.append(indent + entry.name + ("/" if entry.is_dir() else ""))
                if entry.is_dir():
                    lines.extend(_tree(entry, depth + 1))
            return lines

        tree_lines = _tree(root)
        tree_str = f"{root.name}/\n" + "\n".join(tree_lines) if tree_lines else f"{root.name}/  (empty)"

        # ── File count by extension ──
        ext_count: dict[str, int] = {}
        total = 0
        for p in root.rglob("*"):
            if p.is_file() and not any(part in SKIP or part.startswith(".") for part in p.parts):
                ext = p.suffix.lower() or "(no ext)"
                ext_count[ext] = ext_count.get(ext, 0) + 1
                total += 1
        ext_summary = ", ".join(
            f"{ext}: {n}" for ext, n in sorted(ext_count.items(), key=lambda x: -x[1])[:8]
        )

        # ── Key file snapshots ──
        KEY_FILES = ["README.md", "package.json", "requirements.txt", "pyproject.toml",
                     "Makefile", "Dockerfile", "docker-compose.yml", ".env.example"]
        snippets = []
        for name in KEY_FILES:
            kf = root / name
            if kf.exists():
                try:
                    content = kf.read_text(encoding="utf-8", errors="replace")[:400]
                    snippets.append(f"### {name}\n{content.strip()}")
                except Exception:
                    pass

        parts = [
            f"## Project: {root.name}",
            f"Path: {root}",
            f"Total files: {total}  |  By type: {ext_summary or 'n/a'}",
            "",
            "## Directory structure",
            tree_str,
        ]
        if snippets:
            parts += ["", "## Key files"] + snippets

        return "\n".join(parts)
    except Exception as e:
        return f"Error generating project overview: {e}"

@opalatex_tool(
    name="write_content_pos",
    is_safe=False,
    description=(
        "Insert content into an EXISTING file starting at a specific line number (1-indexed). "
        "The new content will be inserted just before the specified line. "
        "The file must already exist: use write_file to create a new file, and "
        "replace_content_range when you need to replace or remove existing lines."
    )
)
def write_content_pos(path: str, content: str, pos: int) -> str:
    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update("write_content_pos", f"path={_preview(resolved)} pos={pos}")
    
    try:
        if os.path.isdir(resolved):
            raise ValueError(f"Error: '{_preview(resolved)}' is a directory. Cannot write to a directory.")
        if not os.path.exists(resolved):
            raise ValueError(
                f"Error: file not found: {_preview(resolved)}. write_content_pos only "
                "inserts into a file that already exists. To create a new file, call "
                "write_file with the full content instead."
            )
    except OSError as e:
        raise ValueError(f"Error: invalid path argument ({e.strerror}).")

    try:
        content = _decode_escape_sequences(content)
        with open(resolved, "r", encoding="utf-8", newline="") as f:
            lines = f.readlines()

        idx = max(0, min(pos - 1, len(lines)))
        
        if content and not content.endswith('\n'):
            content += '\n'
            
        lines.insert(idx, content)
        
        with open(resolved, "w", encoding="utf-8", newline="") as f:
            f.writelines(lines)
        try:
            from .code_index import CODE_INDEX
            CODE_INDEX.rebuild_file(resolved)
        except Exception:
            pass
            
        return f"Successfully inserted content at line {pos} in {_preview(resolved)}."
    except Exception as e:
        raise ValueError(f"Error writing to {_preview(resolved)}: {e}")


@opalatex_tool(
    name="replace_content_range",
    is_safe=False,
    description=(
        "Replace an inclusive 1-indexed line range in an existing file. "
        "Use this for surgical edits to large files instead of rewriting the entire file. "
        "Pass an empty content string to delete the selected lines."
    )
)
def replace_content_range(path: str, start_pos: int, end_pos: int, content: str) -> str:
    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update(
        "replace_content_range",
        f"path={_preview(resolved)} lines={start_pos}-{end_pos}",
    )

    try:
        if os.path.isdir(resolved):
            raise ValueError(f"Error: '{_preview(resolved)}' is a directory. Cannot write to a directory.")
        if not os.path.exists(resolved):
            raise ValueError(
                f"Error: file not found: {_preview(resolved)}. replace_content_range only "
                "edits a file that already exists. To create a new file, call write_file "
                "with the full content instead."
            )
    except OSError as e:
        raise ValueError(f"Error: invalid path argument ({e.strerror}).")

    if start_pos < 1 or end_pos < 1:
        raise ValueError("start_pos and end_pos must be 1-indexed positive integers.")
    if end_pos < start_pos:
        raise ValueError("end_pos must be greater than or equal to start_pos.")

    try:
        with open(resolved, "r", encoding="utf-8", newline="") as f:
            lines = f.readlines()

        if start_pos > len(lines):
            raise ValueError(f"Error: start_pos {start_pos} is beyond the end of the file (total lines: {len(lines)}).")

        start_idx = start_pos - 1
        end_idx = min(end_pos, len(lines))
        replacement = _decode_escape_sequences(content)
        replacement_lines = []
        if replacement:
            if not replacement.endswith("\n"):
                replacement += "\n"
            replacement_lines = replacement.splitlines(keepends=True)

        lines[start_idx:end_idx] = replacement_lines

        with open(resolved, "w", encoding="utf-8", newline="") as f:
            f.writelines(lines)
        try:
            from .code_index import CODE_INDEX
            CODE_INDEX.rebuild_file(resolved)
        except Exception:
            pass

        return (
            f"Successfully replaced lines {start_pos}-{end_pos} "
            f"in {_preview(resolved)}."
        )
    except Exception as e:
        raise ValueError(f"Error replacing content in {_preview(resolved)}: {e}")

@opalatex_tool(
    name="read_content_pos",
    is_safe=True,
    description=(
        "Page through a file by 1-indexed line numbers (inclusive start_pos/end_pos). "
        "This is the paging tool read_file points to when a whole-file read is too big "
        "for the remaining context window: call search_code first to find a target line, "
        "or start at (1, N) to page sequentially from the top. If the requested range is "
        "larger than what still fits the context budget, the returned text is capped and "
        "a trailing note reports the total line count, how many lines remain, and the "
        "start_pos to pass on the next call — always re-read that note rather than "
        "assuming the full requested range came back."
    )
)
def read_content_pos(path: str, start_pos: int, end_pos: int) -> str:
    try:
        resolved = _resolve_path(path)
    except Exception as e:
        raise ValueError(f"Error resolving path: {e}")

    AGENT_PROGRESS.update("read_content_pos", f"path={_preview(resolved)} lines={start_pos}-{end_pos}")

    try:
        if os.path.isdir(resolved):
            raise ValueError(f"Error: '{_preview(resolved)}' is a directory, not a file.")
        if not os.path.exists(resolved):
            raise ValueError(f"Error: file not found: {_preview(resolved)}.")
    except OSError as e:
        raise ValueError(f"Error: invalid path argument ({e.strerror}).")

    if start_pos < 1 or end_pos < 1:
        raise ValueError("start_pos and end_pos must be 1-indexed positive integers.")
    if end_pos < start_pos:
        raise ValueError("end_pos must be greater than or equal to start_pos.")

    try:
        lines = _read_text_file(resolved).splitlines(keepends=True)
        total_lines = len(lines)

        start_idx = start_pos - 1

        if start_idx >= total_lines:
            return (
                f"No content read: start_pos {start_pos} is beyond the end of the file "
                f"(total lines: {total_lines}). Use start_pos <= {total_lines} or locate "
                "the target section before retrying."
            )

        requested_end_idx = min(total_lines, end_pos)

        # A page must fit the remaining context budget on its own: read_content_pos is
        # the escape hatch read_file's own refusal points models toward, so it cannot
        # reproduce the same unbounded-read overflow just because the caller asked for
        # a wide range. Cap the slice instead of refusing outright -- paging tools are
        # expected to return partial pages -- but always say so explicitly (PROJECT_DESIGN
        # 2.6, no silent substitution): a capped page must never look like the full
        # requested range.
        budget_chars = free_context_chars()
        if budget_chars <= 0:
            window = context_window_tokens()
            used = used_context_tokens()
            raise ValueError(
                f"Error: the context window is exhausted ({used:,} of {window:,} tokens "
                f"already used), so '{_preview(resolved)}' cannot be read further. "
                "Summarize what you already know and answer, or ask the user to start a "
                "new chat or raise num_ctx in the project settings."
            )

        end_idx = start_idx
        consumed = 0
        for idx in range(start_idx, requested_end_idx):
            consumed += len(lines[idx])
            if consumed > budget_chars and end_idx > start_idx:
                break
            end_idx = idx + 1
            if consumed > budget_chars:
                break

        selected_lines = lines[start_idx:end_idx]
        content = "".join(selected_lines)
        line_truncated = False
        if len(content) > budget_chars:
            # A single line alone exceeds the whole page budget (e.g. a minified JSON
            # line); still make forward progress instead of returning nothing.
            content = content[:budget_chars]
            line_truncated = True

        actual_end_pos = end_idx
        remaining_lines = total_lines - end_idx

        notes = []
        if line_truncated:
            notes.append(f"line {actual_end_pos} was itself too long to fit and was cut short")
        elif end_idx < requested_end_idx:
            notes.append(f"capped to fit the remaining context budget (you asked through line {end_pos})")
        if remaining_lines > 0:
            notes.append(f"{remaining_lines:,} more line(s) remain")

        if notes:
            content += (
                f"\n\n[Showing lines {start_pos}-{actual_end_pos} of {total_lines:,} total "
                f"lines in '{_preview(resolved)}': " + "; ".join(notes) + ". "
                f"Call read_content_pos({path!r}, {actual_end_pos + 1}, <end_pos>) to continue.]"
            )

        return content
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Error reading {_preview(resolved)}: {e}")

def _rel(path: str, root: str) -> str:
    try:
        return os.path.relpath(path, root)
    except ValueError:
        return path

def get_workspace_action_tools():
    """Tools that act on the workspace: file writes and command execution.

    Single source for what a skill worker can do to the project and for what the
    "direct" orchestrator policy grants (`memgpt_runtime.build_chat_orchestrator`),
    so the two roles cannot drift apart -- the orchestrator used to get the write
    tools listed by hand here and no terminal at all, which left a "direct" model
    able to rewrite `main.tex` but unable to run `pdflatex` on it: half an
    authority, forcing a delegation round-trip mid-task. "delegate" withholds
    exactly this list.

    Every execution tool here is declared `is_safe=False`, so the `opalatex_tool`
    gate still blocks it in plan mode and asks the user in edit mode no matter
    which role calls it.
    """
    return [
        write_file,
        write_content_pos,
        replace_content_range,
        create_docx_file,
        create_pptx_file,
        export_tex_to_docx,
        generate_image,
        run_command,
        run_python_script,
        run_background_command,
        run_interactive_command,
    ]


def get_available_tools():
    """Tools handed to a skill worker (`memgpt_runtime.build_run_skill_tool`)."""
    return [
        search_conversation_history,
        update_achievements_memory,
        get_project_overview,
        search_code,
        read_file,
        read_content_pos,
        analyze_image,
        ask_question,
        get_editor_state,
        *get_workspace_action_tools(),
    ]


# ─── Editor State ─────────────────────────────────────────────────────────────
# The GUI stages what it has open — tabs, focused file, live buffer, selection —
# into <project>/.opalatex/_editor_state.json at the start of every agent turn
# (agent_stdin.handle_run). Reading it back is a pure lookup of data the
# front-end already sent, so the tool is safe: it answers in plan mode without a
# permission prompt, which is exactly when knowing what the user is looking at
# matters most.

EDITOR_STATE_FILENAME = "_editor_state.json"


def editor_state_path(project_path: str) -> str:
    """Absolute path of the staged editor state file for *project_path*."""
    return os.path.join(project_path, ".opalatex", EDITOR_STATE_FILENAME)


def _load_editor_state(project_path: str) -> dict:
    """Return the staged editor state, or {} when nothing has been staged yet."""
    path = editor_state_path(project_path)
    if not os.path.isfile(path):
        return {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            state = json.load(f)
    except (OSError, ValueError) as e:
        raise ValueError(f"Error reading the staged editor state: {e}")
    return state if isinstance(state, dict) else {}


@opalatex_tool(
    name="get_editor_state",
    is_safe=True,
    description=(
        "Report what the user currently has open in the IDE editor: the list of open tabs, "
        "which tab is focused, and the text the user has selected. Call this whenever the user "
        "points at the editor instead of naming a path — 'this file', 'the open document', "
        "'the selected text', 'fix this here' — so you act on what they are looking at instead "
        "of guessing. Pass include_content=True to also get the focused tab's live buffer, which "
        "reflects edits the user has not saved to disk yet; when the saved file is what you need, "
        "read it with read_file instead."
    ),
)
def get_editor_state(include_content: bool = False) -> str:
    AGENT_PROGRESS.update("get_editor_state", f"include_content={include_content}")

    project_path = get_project_path()
    state = _load_editor_state(project_path)
    if not state:
        return (
            "No editor state has been staged for this project yet. The IDE writes it at the "
            "start of each agent turn, so this normally means the request did not come from "
            "the IDE editor. Ask the user which file they mean, or locate it with "
            "get_project_overview / search_code."
        )

    current_file = str(state.get("current_file", "") or "").replace("\\", "/")
    open_files = [str(p).replace("\\", "/") for p in state.get("open_files") or [] if str(p).strip()]
    selected_text = str(state.get("selected_text", "") or "")
    editor_content = str(state.get("editor_content", "") or "")

    lines = ["## Editor State"]

    if open_files:
        listed = "\n".join(
            f"  {i}. `{p}`" + ("  <- focused" if p == current_file else "")
            for i, p in enumerate(open_files, 1)
        )
        lines.append(f"- **Open tabs ({len(open_files)})**:\n{listed}")
    else:
        # Absent (not empty) when the turn came from a surface that does not track
        # tabs, e.g. the inline editor. Say so rather than reporting "no tabs open".
        lines.append("- **Open tabs**: not reported by this turn's caller")

    if current_file:
        lines.append(f"- **Focused file**: `{current_file}`")
        if open_files and current_file not in open_files:
            lines.append(
                "  (the focused file is not among the reported tabs — the user may have "
                "changed tabs after this turn started)"
            )
    else:
        lines.append("- **Focused file**: (no file open in the editor)")

    if selected_text:
        lines.append(
            f"\n### Selected text ({len(selected_text.splitlines())} lines, "
            f"{len(selected_text):,} chars)\n```\n"
            + _truncate_to_context_budget(selected_text, reserve_pct=25)
            + "\n```"
        )
    else:
        lines.append("- **Selection**: (nothing selected)")

    if include_content:
        if editor_content:
            lines.append(
                f"\n### Live buffer of `{current_file or 'the focused tab'}` "
                f"({len(editor_content.splitlines()):,} lines, {len(editor_content):,} chars)\n"
                "This is the editor buffer, which may contain unsaved edits that differ from "
                "the file on disk.\n```\n"
                + _truncate_to_context_budget(editor_content)
                + "\n```"
            )
        else:
            lines.append("\n- **Live buffer**: (empty)")
    elif editor_content:
        lines.append(
            f"- **Live buffer**: {len(editor_content):,} chars available "
            "— call get_editor_state(include_content=True) to read it."
        )

    return "\n".join(lines)


# ─── Achievements Memory ─────────────────────────────────────────────────────
@opalatex_tool(
    name="update_achievements_memory", 
    is_safe=True,
    description=(
        "Update the transient Achievements Memory with a summary of the important milestones and tasks "
        "you have accomplished so far during this turn. Keep it concise (e.g. bullet points). "
        "This memory is passed to worker sub-agents so they know what has already been done. "
        "Use this FREQUENTLY for: discovering an important file, finishing a heartbeat iteration, "
        "successfully reading/writing a file, or finding the root cause of a bug."
    )
)
def update_achievements_memory(summary: str) -> str:
    global TURN_ACHIEVEMENTS
    AGENT_PROGRESS.update("update_achievements_memory", f"summary={summary[:50]}")
    TURN_ACHIEVEMENTS = summary
    try:
        from opalatex.agent_stdin import print_event
        print_event("achievements_update", {"content": summary})
    except Exception as e:
        pass
    return "Achievements Memory updated successfully."

# ─── Long-Term Memory (MemGPT-style) ──────────────────────────────────────────
@opalatex_tool(name="read_core_memory", is_safe=True, description="Read the global 'Core Memory' of the project. Contains rules, persistent context, and architectural decisions you should follow.")
def read_core_memory() -> str:
    AGENT_PROGRESS.update("read_core_memory")
    if not _PROJECT_SESSION:
        return "Core memory not available (no active session)."
    
    if getattr(_PROJECT_SESSION, "use_shared_memory", True):
        return getattr(_PROJECT_SESSION, "core_memory", "") or "(Core memory is empty)"
    else:
        if not _PROJECT_STORE:
            return "(Core memory is empty)"
        mem = _PROJECT_STORE.get_chat_core_memory(_PROJECT_SESSION.name, getattr(_PROJECT_SESSION, "current_chat_id", "main"))
        return mem or "(Core memory is empty)"

@opalatex_tool(name="append_core_memory", is_safe=True, description="Append a new persistent rule, context, or decision to the Core Memory. Do this when you learn something about the user's preferences, demands or the project's state that you want to remember across different executions.")
def append_core_memory(content: str) -> str:
    AGENT_PROGRESS.update("append_core_memory", f"content={_preview(content)}")
    if not _PROJECT_SESSION or not _PROJECT_STORE:
        return "Core memory not available (no active session)."
    
    use_shared = getattr(_PROJECT_SESSION, "use_shared_memory", True)
    
    if use_shared:
        current = getattr(_PROJECT_SESSION, "core_memory", "")
        new_mem = current + ("\n" if current else "") + "- " + content
        _PROJECT_SESSION.core_memory = new_mem
        try:
            _PROJECT_STORE.save(_PROJECT_SESSION)
            return "Successfully appended to global Core Memory."
        except Exception as e:
            raise ValueError(f"Error saving Core Memory: {e}")
    else:
        chat_id = getattr(_PROJECT_SESSION, "current_chat_id", "main")
        current = _PROJECT_STORE.get_chat_core_memory(_PROJECT_SESSION.name, chat_id)
        new_mem = current + ("\n" if current else "") + "- " + content
        try:
            _PROJECT_STORE.update_chat_core_memory(_PROJECT_SESSION.name, chat_id, new_mem)
            return f"Successfully appended to chat '{chat_id}' isolated Core Memory."
        except Exception as e:
            raise ValueError(f"Error saving isolated Core Memory: {e}")

@opalatex_tool(name="search_conversation_history", is_safe=True, description="Search through the past conversations of this project using semantic search (RAG) to remember previous context, decisions, or user instructions. Use this when you need context about past tasks.")
def search_conversation_history(query: str, limit: int = 5) -> str:
    AGENT_PROGRESS.update("search_conversation_history", f"query={_preview(query)}")
    if not _PROJECT_SESSION:
        return "Archival search not available (no active session)."
    
    try:
        from .archival import search_archival
        chat_id = None
        if not getattr(_PROJECT_SESSION, "use_shared_memory", True):
            chat_id = getattr(_PROJECT_SESSION, "current_chat_id", "main")
        results = search_archival(_PROJECT_SESSION.name, query, limit=limit, chat_id=chat_id)
        if not results:
            return f"No results found in archival memory for query: '{query}'"
            
        out = [f"Found {len(results)} results in Archival Memory:"]
        for r in results:
            out.append(f"[{r['timestamp']} | {r['role'].upper()}] {r['content']}")
        return "\n".join(out)
    except Exception as e:
        raise ValueError(f"Error searching Archival Memory: {e}")

def _record_mode_event(content: str) -> None:
    """Record a state-change event as diagnostic activity, not as a chat message.

    ``project_history`` holds only the visible conversation (what the model reads);
    audit entries belong to ``project_activity`` (PROJECT_DESIGN 2.5). The agent
    still learns the outcome within the turn: it is returned in the ``create_plan``
    tool result.

    This is a fire-and-forget helper: it must never raise, because it is called
    inside tool functions and an exception here would surface as a tool error.
    """
    if not _PROJECT_SESSION or not _PROJECT_STORE:
        return
    try:
        _PROJECT_STORE.append_activity(
            _PROJECT_SESSION,
            "plan_decision",
            content=content,
            agent="chat_orchestrator",
        )
    except Exception:
        pass  # never let history recording crash the tool


@opalatex_tool(
    name="create_plan",
    is_safe=True,
    description=(
        "Presents a proposed implementation plan to the user for approval. "
        "Use this tool when you are in 'plan' mode and have gathered enough context to propose a plan. "
        "This tool halts execution and waits for the user to approve or reject the plan."
    )
)
async def create_plan(plan_content: str) -> str:
    AGENT_PROGRESS.update("create_plan")
    if not _PROJECT_SESSION:
        return "Failed: No active session."
    
    T.info("Presenting Proposed Plan to user for approval...")
    
    # Emit custom input_request with markdown_content instead of using T.aconfirm
    import uuid
    import asyncio
    from opalatex.agent_stdin import print_event, _gui_input_pending
    
    loop = asyncio.get_event_loop()
    req_id = str(uuid.uuid4())
    fut = loop.create_future()
    _gui_input_pending[req_id] = fut
    
    print_event("input_request", {
        "id": req_id,
        "prompt": "Review the proposed plan below:",
        "markdown_content": plan_content,
        "type": "confirm",
        "options": ["yes", "no"],
        "default": "yes"
    })
    
    try:
        raw = await asyncio.wait_for(asyncio.shield(fut), timeout=86400.0) # wait up to 24h
        
        # Try parsing JSON if the frontend sent the edited content
        import json
        try:
            data = json.loads(raw)
            approved_str = data.get("response", "no")
            edited_plan = data.get("editedContent", plan_content)
        except Exception:
            approved_str = raw.strip()
            edited_plan = plan_content
            
        approved = approved_str.lower() in ("yes", "y", "s", "sim", "true", "1")
    except asyncio.TimeoutError:
        approved = True
        edited_plan = plan_content
    finally:
        _gui_input_pending.pop(req_id, None)
    
    if approved:
        prev_mode = getattr(_PROJECT_SESSION, "_initial_mode", "plan")
        _PROJECT_SESSION.mode = "auto"
        if edited_plan != plan_content:
            _PROJECT_SESSION.plan_text = edited_plan
        # Record approval immediately in chat history. Do not save the whole
        # project here, because the mode is only temporarily auto for this turn.
        _record_mode_event(
            f"[PLAN APPROVED] The user approved the proposed plan. "
            f"Temporary mode for this turn: '{prev_mode}' -> 'auto'. "
            "The project mode must be restored at the end of the turn."
        )
        return f"The user APPROVED the plan. The system is temporarily in 'auto' mode for this turn only. Proceed to execute the plan.\n\nPlan Content:\n{edited_plan}"
    else:
        # Record rejection immediately in chat history
        _record_mode_event(
            "[PLAN REJECTED] The user rejected the proposed plan. "
            "The agent must NOT proceed with execution. "
            "Wait for user feedback before proposing a new plan."
        )
        # Throw an error to halt the agent immediately and prevent it from continuing
        raise ValueError("The user REJECTED the plan. Wait for the user to provide feedback in the chat.")


def context_window_tokens() -> int:
    """Return the active context window, which is what num_ctx actually caps.

    Resolved via config.resolve_effective_num_ctx: an explicit project
    override wins, otherwise the model's catalog entry, otherwise the
    local/cloud heuristic default.
    """
    from .config import resolve_effective_num_ctx
    try:
        return resolve_effective_num_ctx("memgpt")
    except Exception:
        return 8192


def used_context_tokens() -> int:
    """Return how much of the window is already occupied.

    Prefers the provider-reported `prompt_tokens` of the last call (see
    opalatex.token_usage), which is the only figure that counts the system
    prompt, the tool schemas and previous tool results. The char/4 estimate over
    project_history is the fallback before the first measured call, and it
    understates the real occupancy badly.
    """
    from opalatex.token_usage import get_context_prompt_tokens

    measured = get_context_prompt_tokens()
    if measured is not None:
        return measured
    history = getattr(_PROJECT_SESSION, "history", None) or []
    return len(json.dumps(history)) // 4


def free_context_chars(reserve_pct: int = 50) -> int:
    """Return how many characters a tool result may still add.

    Only *reserve_pct* of the remaining window is granted to any single tool
    result, leaving room for the rest of the turn and the model's own reply.
    """
    free_tokens = max(0, context_window_tokens() - used_context_tokens())
    return int(free_tokens * 4 * reserve_pct / 100)


def _truncate_to_context_budget(text: str, reserve_pct: int = 50) -> str:
    """Truncate *text* so it fits within the active model's free context budget."""
    allowed_chars = free_context_chars(reserve_pct)
    if allowed_chars <= 0 or len(text) <= allowed_chars:
        return text
    return (
        text[:allowed_chars]
        + f"\n\n[Result truncated: {len(text):,} chars total, "
        f"{allowed_chars:,} shown ({reserve_pct}% of free context budget)]"
    )


@opalatex_tool(
    name="web_search",
    is_safe=True,
    description=(
        "Search the web for current information, documentation, news, or any topic. "
        "Use this when you need information that may be recent, external, or not in your training data. "
        "Returns a list of results with titles, URLs and snippets. "
        "Requires web search to be enabled in OpalaTex settings."
    ),
)
def web_search(query: str, max_results: int = 5) -> str:
    """Search the web using DuckDuckGo (default) or the configured MCP server."""
    AGENT_PROGRESS.update("web_search", f"query={_preview(query)}")
    try:
        from .web_search_config import is_enabled, do_search
    except ImportError as exc:
        return f"[web_search] Configuration module unavailable: {exc}"

    if not is_enabled():
        return (
            "[web_search] Web search is currently disabled. "
            "Enable it via the Web Search toggle in the OpalaTex chat panel."
        )

    # Bridge the async do_search into the synchronous tool interface.
    # Always use a fresh ThreadPoolExecutor so that asyncio.run() creates its own
    # event loop in a clean thread — this is safe regardless of whether the caller
    # is the main thread, an async worker thread (agenticblocks), or any other
    # context.  Using asyncio.get_event_loop() in Python 3.10+ raises RuntimeError
    # when called from a thread that never had an event loop (e.g. 'asyncio_0').
    import concurrent.futures
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(__import__("asyncio").run, do_search(query, max_results))
            result = future.result(timeout=30)
    except Exception as exc:
        return f"[web_search] Search failed: {exc}"

    return _truncate_to_context_budget(result)
